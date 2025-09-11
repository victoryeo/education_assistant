import asyncio
import json
import os
import sys

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from mcp.client.stdio import stdio_client
from mcp.client.session import ClientSession
from typing import Any, Dict
from mcp import ClientSession, StdioServerParameters
from pptx import Presentation
from pptx.util import Inches, Pt

# Server configuration class
class ServerConfig:
    def __init__(self, command: str, args: list[str]):
        self.command = command
        self.args = args
        self.env = os.environ.copy()
        self.cwd = os.getcwd()
        # Add missing attributes expected by MCP client
        self.encoding = "utf-8"
        self.encoding_error_handler = "strict"
        self.stderr = None  # Let subprocess handle stderr normally
        self.capabilities = None

async def client_pdf_generate():
    """Example client usage demonstrating basic MCP operations"""
    
    print("=== MCP Client Usage Example ===\n")
    # Get the absolute path to the server script
    server_script = os.path.abspath("multiagent_mcp_server.py")
    
    # Create server config
    server_config = ServerConfig(
        command=sys.executable,
        args=[server_script]
    )
    async with stdio_client(server_config) as (read, write):
        async with ClientSession(read, write) as session:
            # Initialize session
            await session.initialize()
            print("✅ Session initialized")
            
            # List available tools
            tools = await session.list_tools()
            print(f"📋 Available tools: {[tool.name for tool in tools.tools]}")
            
            # List available resources
            resources = await session.list_resources()
            print(f"📚 Available resources: {[res.name for res in resources.resources]}")
            
            # Example 1: Process a learning request
            print("\n--- Example 1: Process Learning Request ---")
            result = await session.call_tool(
                "process_message",
                {
                    "user_input": "I want to learn about photosynthesis",
                    "user_id": "demo_user",
                    "category": "biology"
                }
            )
            # Handle the response content
            if not result.content:
                print("Error: Empty response received from server")
                return
                
            response_text = result.content[0].text
            print(f"\n--- Raw Response ---\n{response_text}\n--- End Raw Response ---\n")
            
            # Check for validation errors in the response
            if "validation error" in response_text.lower() or "validation_error" in response_text.lower():
                print("⚠️ Validation error in response:")
                print(response_text)
                return
                
            try:
                # Try to parse as JSON
                agent_response = None
                response_data = json.loads(response_text)
                print(f"Response type: {type(response_data).__name__}")
                if isinstance(response_data, dict):
                    if 'response' in response_data:
                        print(f"✅ Agent response: {response_data['response'][:200]}...")
                        agent_response = response_data['response']
                    elif 'message' in response_data:
                        print(f"✅ Message: {response_data['message']}")
                        agent_response = response_data['message']
                    elif 'content' in response_data:
                        print(f"✅ Content: {response_data['content']}")
                        response_list = response_data['content']
                        response_dict = response_list[0]
                        # The 'text' key holds the JSON string
                        response_text_json = response_dict.get('text', '{}')
                        agent_response = json.dumps(response_text_json)
                    else:
                        print(f"✅ Response data: {json.dumps(response_data, indent=2)}")
                        agent_response = json.dumps(response_data, indent=2)

                    # genetate PDF file
                    print(f"Agent Response type: {type(agent_response).__name__}")
                    print(agent_response)
                    pdf_filename = "learning_report.pdf"
                    doc = SimpleDocTemplate(pdf_filename, pagesize=letter)
                    styles = getSampleStyleSheet()
                    
                    # Define a story (list of elements) to be added to the PDF
                    story = []
                    
                    # Add a title and user query
                    story.append(Paragraph("Learning Report: Photosynthesis", styles['Title']))
                    story.append(Spacer(1, 0.2 * inch))
                    story.append(Paragraph("<b>User Query:</b> I want to learn about photosynthesis", styles['Normal']))
                    story.append(Spacer(1, 0.2 * inch))
                    
                    # Add the agent's response to the story
                    if agent_response:
                        story.append(Paragraph(agent_response, styles['BodyText']))
                    
                    # Build the PDF file
                    doc.build(story)
                    print(f"\n✅ PDF file '{pdf_filename}' generated successfully!")

                    # generate PPT file
                    prs = Presentation()
                    slide_layout = prs.slide_layouts[5]  # Using blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = "Learning Report: Photosynthesis"
                    
                    # Add content
                    left = width = height = Inches(1)  # 1 inch margin
                    top = Inches(1.5)
                    width = Inches(8.5)  # Slide width minus margins
                    height = Inches(5)    # Reasonable height for content
                    
                    # Add text box for the content
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    
                    # Add user query
                    p = tf.add_paragraph()
                    p.text = "User Query: I want to learn about photosynthesis"
                    p.font.bold = True
                    
                    # Add agent response
                    if agent_response:
                        p = tf.add_paragraph()
                        p.text = agent_response if isinstance(agent_response, str) else str(agent_response)
                        p.space_after = Pt(12)  # Add some space after the paragraph
                    
                    # Save the presentation
                    pptx_filename = "learning_report.pptx"
                    prs.save(pptx_filename)
                    print(f"✅ PowerPoint file '{pptx_filename}' generated successfully!")
                    
                else:
                    print(f"✅ Response: {response_data}")
            except json.JSONDecodeError:
                print(f"📝 Raw response (non-JSON): {response_text[:500]}...")
                                
            except Exception as e:
                print(f"❌ Error reading agent: {str(e)}")
            
            print("\n✅ Example usage completed successfully!")

# Main execution
if __name__ == "__main__":
    asyncio.run(client_pdf_generate())
          
    print("\n=== Usage Instructions ===")
    print("1. Set up your .env file with required API keys")
    print("2. Install requirements: pip install -r requirements.txt") 
    print("3. Run server: python multiagent_mcp_server.py")
    print("4. Run tests: python mcp_pdf.py")